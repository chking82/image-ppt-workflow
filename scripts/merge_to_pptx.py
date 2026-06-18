#!/usr/bin/env python3
"""
merge_to_pptx.py - 合并幻灯片图片 + 演讲备注为 PPTX
所有文字必须在生图时嵌入到图片中，本脚本只做两件事：
  1. 把图片铺满每一页幻灯片
  2. 把演讲备注写入 PowerPoint 备注区

用法:
  python3 merge_to_pptx.py \\
    --slides <项目>/slides/ \\
    --notes <项目>/speaker-notes.md \\
    --output <项目>/<主题名称>.pptx
"""

import argparse
import glob
import os
import re
import sys
from pptx import Presentation
from pptx.util import Inches, Emu

A4_WIDTH = Emu(12192000)
A4_HEIGHT = Emu(6858000)


def parse_speaker_notes(notes_path):
    """解析 speaker-notes.md，返回按页索引的口述稿列表"""
    if not notes_path or not os.path.isfile(notes_path):
        return []

    with open(notes_path, "r", encoding="utf-8") as f:
        content = f.read()

    # 按 "## Slide N" 分割，提取每页备注
    notes = {}
    for match in re.finditer(r"## Slide (\d+)[^\n]*\n\n(.*?)(?=## Slide |\Z)", content, re.DOTALL):
        page_num = int(match.group(1))
        note_text = match.group(2).strip()
        notes[page_num] = note_text

    # 返回按页码排序的备注列表
    max_page = max(notes.keys()) if notes else 0
    return [notes.get(i, "") for i in range(1, max_page + 1)]


def create_pptx(slide_images, speaker_notes, output_path):
    """创建 PPTX 文件"""
    prs = Presentation()
    prs.slide_width = A4_WIDTH
    prs.slide_height = A4_HEIGHT
    blank_layout = prs.slide_layouts[6]  # 空白版式

    for idx, img_path in enumerate(slide_images):
        slide = prs.slides.add_slide(blank_layout)

        # 添加图片到幻灯片
        slide.shapes.add_picture(img_path, 0, 0, A4_WIDTH, A4_HEIGHT)

        # 添加演讲备注
        if idx < len(speaker_notes) and speaker_notes[idx]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes[idx]

    prs.save(output_path)
    print(f"✅ PPTX 已保存: {output_path} ({len(slide_images)} 页)")


def main():
    parser = argparse.ArgumentParser(description="合并幻灯片图片 + 演讲备注为 PPTX")
    parser.add_argument("--slides", required=True, help="幻灯片图片目录")
    parser.add_argument("--notes", required=True, help="speaker-notes.md 路径")
    parser.add_argument("--output", required=True, help="输出 PPTX 文件路径")
    parser.add_argument("--require-manifest", help="manifest.json 路径 (传则校验, 缺页则拒绝合成)")
    parser.add_argument("--prompts-dir", help="prompts/ 目录 (跟 manifest 对比用)")
    args = parser.parse_args()

    from pathlib import Path as _P
    import subprocess as _sp

    if _P(args.output).name == "output.pptx" and os.environ.get("IMAGE_PPT_ALLOW_GENERIC_OUTPUT") != "1":
        print("❌ 输出文件名不能使用 output.pptx，请使用基于主题的文件名", file=sys.stderr)
        sys.exit(1)

    # manifest 硬门禁：正式流程必须提供；如需旧项目调试，显式设置 IMAGE_PPT_ALLOW_NO_MANIFEST=1。
    if not args.require_manifest and os.environ.get("IMAGE_PPT_ALLOW_NO_MANIFEST") != "1":
        print("❌ 正式流程必须传 --require-manifest slides-manifest.json", file=sys.stderr)
        sys.exit(1)

    if args.require_manifest:
        manifest_path = _P(args.require_manifest)
        if not manifest_path.exists():
            print(f"❌ manifest 不存在: {manifest_path}", file=sys.stderr)
            sys.exit(1)
        manifest_py = _P(__file__).parent / "manifest.py"
        checks = [
            [sys.executable, str(manifest_py), "check", str(manifest_path)],
            [sys.executable, str(manifest_py), "reconcile", str(manifest_path), args.slides],
        ]
        prompts_dir = args.prompts_dir or "prompts"
        if _P(prompts_dir).exists():
            checks.append([sys.executable, str(manifest_py), "reconcile-prompts", str(manifest_path), prompts_dir])
        for cmd in checks:
            result = _sp.run(cmd, capture_output=True, text=True)
            if result.returncode != 0:
                print("❌ manifest 门禁失败:", " ".join(cmd), file=sys.stderr)
                if result.stdout:
                    print(result.stdout, file=sys.stderr)
                if result.stderr:
                    print(result.stderr, file=sys.stderr)
                sys.exit(1)
        print(f"✅ manifest 门禁通过: {manifest_path}")

    # 获取所有图片，按文件名排序
    image_extensions = ("*.png", "*.jpg", "*.jpeg", "*.bmp", "*.gif")
    images = []
    for ext in image_extensions:
        images.extend(glob.glob(os.path.join(args.slides, ext)))
    images.sort()

    if not images:
        print(f"❌ 在 {args.slides} 中没有找到图片")
        return

    print(f"📷 找到 {len(images)} 张图片")

    # 解析演讲备注
    notes = parse_speaker_notes(args.notes)
    if notes:
        print(f"📝 解析到 {len(notes)} 页备注")
    else:
        print("⚠️  备注文件为空或不存在")

    # 创建 PPTX
    create_pptx(images, notes, args.output)


if __name__ == "__main__":
    main()
